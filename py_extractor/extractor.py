import os, zipfile, shutil
from pydub import AudioSegment
import speech_recognition as sr
from pptx import Presentation
from pathlib import Path


audio_extensions = ('.aiff', '.au', '.mid', '.midi', '.mp3', '.m4a', '.mp4', '.wav', '.wma')
init_dir = os.getcwd()
init_files = os.listdir()
results_dir =  "../results"
ppt_dir = init_dir + '/ppt'
media_dir = ppt_dir + '/media'
output_file_prefix = "converted_"
result_dir_contents = os.listdir(results_dir)

def prepare_voice_file(path):
    audio_file = AudioSegment.from_file(path, format="m4a")
    wav_file = os.path.splitext(path)[0] + '.wav'
    audio_file.export(wav_file, format='wav')
    return wav_file

def extract_audio_files(directory):
    for power in directory:

        fbase = "".join(power.split(".")[:-1]).replace(' ', '_')

        if not power.endswith('.zip') and not power.endswith('.pptx') or power == 'base_library.zip':
            continue
        elif fbase in result_dir_contents:
            print(f"{fbase} already processed! Skipping...")
            continue

        base = os.path.splitext(power)[0]
        power_copy_copy = base + '_copy.pptx'
        new_zip = base + '.zip'
        shutil.copy(power, power_copy_copy)
        os.rename(power_copy_copy , new_zip)
        
        with zipfile.ZipFile(new_zip) as myzip:
            if myzip.testzip() != None:
                print("Some of your media files are corrupted")

            else :
                for file in myzip.namelist() :
                    if file.endswith(audio_extensions):
                        myzip.extract(file)
                
        myzip.close()
        os.remove(new_zip)
        break

    return power

def transcribe(audio_file_paths, result_subdir, pptx_file_name):
    r = sr.Recognizer()
    total_num = len(audio_file_paths)
    processing_count = 1
    pptx_file_path = f"{result_subdir}/{pptx_file_name}"
    for f in audio_file_paths:
        f_name = f.split("/")[-1].split(".")[0]
        slide_num = int(f_name.split('media')[1])
        print(f"Processing {processing_count}/{total_num}...")
        wav_file = prepare_voice_file(f)
        with sr.AudioFile(wav_file) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
            if text:
                with open(f"{result_subdir}/{output_file_prefix}{f_name}.txt", 'w') as out_file:
                    # to_notes(text, pptx_file_path, slide_num - 1)     # write the transcript to the speaker notes
                    out_file.write(text)    # write the transcript to a separate txt
        processing_count += 1

def to_notes(note, pptx_fname, slide_num):
    pptx_file = Path(pptx_fname)
    prs = Presentation(pptx_file)
    slide = prs.slides[slide_num]
    slide.notes_slide.notes_text_frame.text = note
    prs.save(pptx_file)

def main():
    print('Extracting...')

    if os.path.isdir(ppt_dir):
        shutil.rmtree(ppt_dir)
    
    file = extract_audio_files(init_files)

    audio_files = os.listdir(media_dir)
    audio_file_paths = [f"{media_dir}/{f}" for f in audio_files]
    sorted_audio_files = sorted(audio_file_paths, key=lambda x: int(x.split('/')[-1].split('.')[0].split('media')[1]))

    if not sorted_audio_files:
        return

    result_dir_relative = "".join(file.split(".")[:-1]).replace(" ", "_")
    result_subdir = f"{results_dir}/{result_dir_relative}"

    try:
        os.mkdir(result_subdir)
    except: pass

    shutil.copy(file, result_subdir)
    transcribe(sorted_audio_files, result_subdir, file)
    
    shutil.rmtree(ppt_dir)

if __name__ == "__main__":
    main()
